"""Build KiX 7-slide investor PPTX (v2 — benchmark insights edition).

Dark navy background, KiX green accents, 16:9 widescreen, editable in PowerPoint/Keynote.

7-slide structure:
  1. Killer Line                — "$50M vs $49/month"
  2. The 5-Step Funnel          — Game→Reward→Register→Redeem→Return
  3. The 9 Real Competitors     — Direct / Operators / Adjacent
  4. KiX Differentiation        — Blue ocean matrix
  5. System Architecture        — 9-card stack (from v1 slide 2)
  6. 90-Day Bedok Land Grab     — Day 0-30-60-90 timeline
  7. Crystal Clear Call         — 3 lines + founder mantra
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

# ───── colors (match HTML deck) ─────
NAVY = RGBColor(0x0F, 0x17, 0x2A)
NAVY_LIGHT = RGBColor(0x1E, 0x29, 0x3B)
NAVY_BORDER = RGBColor(0x33, 0x41, 0x55)
WHITE = RGBColor(0xF8, 0xFA, 0xFC)
GRAY_LIGHT = RGBColor(0xCB, 0xD5, 0xE1)
GRAY = RGBColor(0x94, 0xA3, 0xB8)
GRAY_DARK = RGBColor(0x64, 0x74, 0x8B)
KIX_GREEN = RGBColor(0x00, 0xB3, 0x41)
KIX_GREEN_BG = RGBColor(0x00, 0x3C, 0x18)
BLUE = RGBColor(0x1A, 0x73, 0xE8)
BLUE_BG = RGBColor(0x0B, 0x21, 0x45)
SLATE_FOOTER = RGBColor(0x47, 0x55, 0x69)
RED_MUTED = RGBColor(0xDC, 0x4C, 0x4C)
AMBER = RGBColor(0xF5, 0xB7, 0x3F)

SCREENSHOT_DIR = Path("/Users/mozat/a-docs/slide-screenshots")
OUT_PATH = Path("/Users/mozat/a-docs/kix-investor-deck.pptx")
TOTAL_SLIDES = 7


def fill_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill_color, line_color=None, line_width=0):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_round_rect(slide, x, y, w, h, fill_color, line_color=None, line_width=0, radius=0.08):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.adjustments[0] = radius
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text(slide, x, y, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT, font="Inter"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    return tb, tf, p


def add_pill(slide, x, y, text, fg=KIX_GREEN, bg=KIX_GREEN_BG, width=2.4):
    pill_w = Inches(width); pill_h = Inches(0.32)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, pill_w, pill_h)
    shape.adjustments[0] = 1.0  # max rounding
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.fill.background()
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    p.text = text
    run = p.runs[0]
    run.font.name = "Inter"; run.font.size = Pt(9); run.font.bold = True
    run.font.color.rgb = fg
    return shape


def add_stat_card(slide, x, y, w, h, num, label, detail):
    add_rect(slide, x, y, w, h, NAVY_LIGHT, NAVY_BORDER, 0.75)
    pad = Inches(0.2)
    add_text(slide, x + pad, y + pad, w - pad*2, Inches(0.6),
             num, 28, KIX_GREEN, bold=True)
    add_text(slide, x + pad, y + Inches(0.85), w - pad*2, Inches(0.25),
             label.upper(), 8, GRAY, bold=True)
    add_text(slide, x + pad, y + Inches(1.15), w - pad*2, Inches(0.75),
             detail, 8, GRAY_LIGHT)


def add_arch_card(slide, x, y, w, h, head, body):
    add_rect(slide, x, y, w, h, NAVY_LIGHT, NAVY_BORDER, 0.5)
    pad = Inches(0.15)
    add_text(slide, x + pad, y + pad, w - pad*2, Inches(0.25),
             head, 8, KIX_GREEN, bold=True)
    add_text(slide, x + pad, y + Inches(0.45), w - pad*2, h - Inches(0.5),
             body, 8, GRAY_LIGHT)


def add_logo(slide, x, y):
    tb = slide.shapes.add_textbox(x, y, Inches(1.5), Inches(0.5))
    tf = tb.text_frame
    tf.margin_left = Emu(0); tf.margin_top = Emu(0)
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = "Ki"
    r1.font.name = "Inter"; r1.font.size = Pt(20); r1.font.bold = True
    r1.font.color.rgb = WHITE
    r2 = p.add_run(); r2.text = "X"
    r2.font.name = "Inter"; r2.font.size = Pt(20); r2.font.bold = True
    r2.font.color.rgb = KIX_GREEN


def add_footer(slide, text):
    add_text(slide, Inches(0.6), Inches(7.05), Inches(12), Inches(0.3),
             text, 8, SLATE_FOOTER)


def add_slide_header(slide, idx, pill_text, headline, subhead=None,
                     headline_size=30, pill_color=KIX_GREEN, pill_bg=KIX_GREEN_BG):
    fill_bg(slide, NAVY)
    add_logo(slide, Inches(0.6), Inches(0.4))
    add_text(slide, Inches(0.6), Inches(0.95), Inches(3), Inches(0.25),
             f"SLIDE {idx} / {TOTAL_SLIDES}", 8, GRAY_DARK, bold=True)
    add_pill(slide, Inches(0.6), Inches(1.3), pill_text, fg=pill_color, bg=pill_bg, width=3.0)
    add_text(slide, Inches(0.6), Inches(1.75), Inches(12.2), Inches(1.0),
             headline, headline_size, WHITE, bold=True)
    if subhead:
        add_text(slide, Inches(0.6), Inches(2.55), Inches(12.2), Inches(0.5),
                 subhead, 14, GRAY)


# ───── build presentation ─────
prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9 widescreen
prs.slide_height = Inches(7.5)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height
blank = prs.slide_layouts[6]


# ═══════════════════════════════════════════════════════════════════
# SLIDE 1 — The Killer Line
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
fill_bg(slide, NAVY)

add_logo(slide, Inches(0.6), Inches(0.4))
add_text(slide, Inches(0.6), Inches(0.95), Inches(3), Inches(0.25),
         f"SLIDE 1 / {TOTAL_SLIDES}", 8, GRAY_DARK, bold=True)
add_pill(slide, Inches(0.6), Inches(1.3), "THE KILLER LINE", width=2.4)

# Big two-tone headline
tb = slide.shapes.add_textbox(Inches(0.6), Inches(1.75), Inches(12.2), Inches(1.4))
tf = tb.text_frame
tf.margin_left = Emu(0); tf.margin_right = Emu(0)
tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
for txt, col, bold in [
    ("What McDonald's pays ", WHITE, True),
    ("$50M", KIX_GREEN, True),
    (" for,\nevery merchant gets for ", WHITE, True),
    ("$49/month", KIX_GREEN, True),
    (".", WHITE, True),
]:
    parts = txt.split("\n")
    for i, part in enumerate(parts):
        if i > 0:
            p = tf.add_paragraph(); p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = part
        r.font.name = "Inter"; r.font.size = Pt(34); r.font.bold = bold
        r.font.color.rgb = col

add_text(slide, Inches(0.6), Inches(3.55), Inches(12.2), Inches(0.5),
         "Self-serve McDonald's Monopoly for every offline merchant.",
         16, GRAY_LIGHT, bold=False)

# 4 stat cards row
card_w = Inches(2.95); card_h = Inches(1.7); card_y = Inches(4.25); gap = Inches(0.15)
for i, (n, l, d) in enumerate([
    ("285K", "Lines of code", "85 routers · 60+ DB models · 6 migrations"),
    ("1248", "Tests passing", "100% pass rate · 0 known failures"),
    ("11", "Locales live", "EN-SG · ZH · 4 SEA · 3 RTL · CLDR plurals"),
    ("9", "Regions ready", "SG · ID · TH · VN · PH · EU · US · BR · IN"),
]):
    add_stat_card(slide, Inches(0.6) + i * (card_w + gap), card_y, card_w, card_h, n, l, d)

# Storefront screenshot strip (full-width below cards)
ss_path = SCREENSHOT_DIR / "storefront.png"
if ss_path.exists():
    box_x = Inches(0.6); box_y = Inches(6.1); ss_w = Inches(12.13); ss_h = Inches(0.85)
    add_rect(slide, box_x, box_y, ss_w, ss_h, NAVY_LIGHT, NAVY_BORDER, 0.5)
    # picture aspect = 16:9-ish wide; keep as banner
    slide.shapes.add_picture(str(ss_path), box_x + Inches(0.05), box_y + Inches(0.05),
                             width=ss_w - Inches(0.1), height=ss_h - Inches(0.1))

add_footer(slide, "KiX Platform · The productized McDonald's Monopoly · michael@mozat.com")


# ═══════════════════════════════════════════════════════════════════
# SLIDE 2 — The 5-Step Funnel
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
add_slide_header(slide, 2, "THE 5-STEP FUNNEL",
                 "Game → Reward → Register → Redeem → Return.",
                 "One stack that owns the whole consumer journey, end-to-end.",
                 headline_size=30)

# Horizontal funnel: 5 chevron-like cards + arrows
steps = [
    ("🎮", "GAME",     "Spin · Match · Scratch"),
    ("🎁", "REWARD",   "Voucher · Credit · NFT"),
    ("📝", "REGISTER", "Phone · OAuth · KYC"),
    ("💳", "REDEEM",   "QR · Wallet · POS"),
    ("🔁", "RETURN",   "Push · Streak · Loyalty"),
]
step_y = Inches(3.25)
step_h = Inches(2.05)
total_w = Inches(12.13)
n = len(steps)
# 5 steps, 4 arrows; arrow ~0.45in
arrow_w = Inches(0.4)
step_w_emu = (total_w - arrow_w * (n - 1)) / n
step_w = Emu(int(step_w_emu))
x = Inches(0.6)
for i, (emoji, label, sub) in enumerate(steps):
    add_round_rect(slide, x, step_y, step_w, step_h, NAVY_LIGHT, KIX_GREEN, 1.0, radius=0.1)
    # emoji
    add_text(slide, x, step_y + Inches(0.25), step_w, Inches(0.7),
             emoji, 40, WHITE, bold=False, align=PP_ALIGN.CENTER)
    # label
    add_text(slide, x, step_y + Inches(1.05), step_w, Inches(0.35),
             label, 16, KIX_GREEN, bold=True, align=PP_ALIGN.CENTER)
    # subtext
    add_text(slide, x, step_y + Inches(1.45), step_w, Inches(0.5),
             sub, 10, GRAY_LIGHT, align=PP_ALIGN.CENTER)
    if i < n - 1:
        ax = x + step_w
        # chevron arrow shape
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax, step_y + Inches(0.85), arrow_w, Inches(0.4))
        arrow.fill.solid(); arrow.fill.fore_color.rgb = KIX_GREEN
        arrow.line.fill.background()
        arrow.shadow.inherit = False
    x = x + step_w + arrow_w

# Insight blocks below
add_text(slide, Inches(0.6), Inches(5.55), Inches(12.2), Inches(0.4),
         "Big brands pay $50M for specialists at each step. We productized the entire stack.",
         14, WHITE, bold=True)

# Key insight callout
ki_x = Inches(0.6); ki_y = Inches(6.15); ki_w = Inches(12.13); ki_h = Inches(0.8)
add_rect(slide, ki_x, ki_y, Inches(0.05), ki_h, KIX_GREEN)
add_rect(slide, ki_x + Inches(0.05), ki_y, ki_w - Inches(0.05), ki_h, NAVY_LIGHT)
add_text(slide, ki_x + Inches(0.25), ki_y + Inches(0.1), ki_w - Inches(0.4), Inches(0.3),
         "KEY INSIGHT", 9, KIX_GREEN, bold=True)
add_text(slide, ki_x + Inches(0.25), ki_y + Inches(0.38), ki_w - Inches(0.4), Inches(0.4),
         "Gamified promotion infrastructure — not loyalty SaaS, not ad-tech. A new category.",
         12, WHITE, bold=True)

add_footer(slide, "5-Step Funnel · The complete consumer journey — owned end-to-end by one stack")


# ═══════════════════════════════════════════════════════════════════
# SLIDE 3 — The 9 Real Competitors
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
add_slide_header(slide, 3, "THE COMPETITIVE LANDSCAPE",
                 "9 real competitors. Zero serve SMBs.",
                 "Benchmark study (May 2026) — every player in gamified promotion, segmented by tier.",
                 headline_size=28)

# 3 tier columns
tier_y = Inches(3.25); tier_h = Inches(3.4); col_gap = Inches(0.25)
total_w = Inches(12.13)
col_w = Emu(int((total_w - col_gap * 2) / 3))

tiers = [
    ("TIER A — DIRECT",
     "Gamified-promotion specialists",
     KIX_GREEN,
     ["Gamify.com", "BRAME", "Playable (Riddle)", "CataBoom"]),
    ("TIER B — OPERATORS",
     "Big-brand agencies & ops",
     AMBER,
     ["tms (McDonald's Monopoly)", "Merkle", "Flarie"]),
    ("TIER C — ADJACENT",
     "Loyalty / engagement plays",
     BLUE,
     ["Realtime Media", "Bunchball (Hive)"]),
]
x = Inches(0.6)
for tname, sub, color, members in tiers:
    add_round_rect(slide, x, tier_y, col_w, tier_h, NAVY_LIGHT, color, 1.25, radius=0.05)
    # Top color bar
    add_rect(slide, x, tier_y, col_w, Inches(0.08), color)
    pad = Inches(0.2)
    add_text(slide, x + pad, tier_y + Inches(0.2), col_w - pad*2, Inches(0.35),
             tname, 11, color, bold=True)
    add_text(slide, x + pad, tier_y + Inches(0.55), col_w - pad*2, Inches(0.3),
             sub, 9, GRAY, bold=False)
    # Divider line
    add_rect(slide, x + pad, tier_y + Inches(0.9), col_w - pad*2, Emu(9525), NAVY_BORDER)
    # Members list
    item_y = tier_y + Inches(1.05)
    for m in members:
        add_text(slide, x + pad, item_y, Inches(0.25), Inches(0.3),
                 "▸", 12, color, bold=True)
        add_text(slide, x + pad + Inches(0.3), item_y + Inches(0.02), col_w - pad*2 - Inches(0.3), Inches(0.3),
                 m, 11, WHITE, bold=False)
        item_y += Inches(0.42)
    x += col_w + col_gap

# Caption callout
cap_y = Inches(6.85); cap_h = Inches(0.4)
add_text(slide, Inches(0.6), cap_y, Inches(12.13), cap_h,
         "None of them can profitably serve SMB merchants at $49/month.",
         13, KIX_GREEN, bold=True, align=PP_ALIGN.CENTER)

add_footer(slide, "Competitive Landscape · 9 benchmarked players · KiX = only self-serve SMB option")


# ═══════════════════════════════════════════════════════════════════
# SLIDE 4 — KiX Differentiation (Blue Ocean Matrix)
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
add_slide_header(slide, 4, "BLUE OCEAN — THE DIFFERENTIATION MATRIX",
                 "Where the 9 competitors stop, KiX begins.",
                 "Self-serve · marginal cost ≈ 0 · SEA-native · cross-brand · 5 minutes to launch.",
                 headline_size=26)

# Matrix table: 9 rows × (label + 3 cols)
mat_x = Inches(0.6); mat_y = Inches(3.2)
mat_w = Inches(12.13); row_h = Inches(0.38)
label_w = Inches(3.3)
col_w_m = Emu(int((mat_w - label_w) / 3))

# Header
add_rect(slide, mat_x, mat_y, mat_w, row_h, NAVY_BORDER)
add_text(slide, mat_x + Inches(0.15), mat_y + Inches(0.07), label_w - Inches(0.15), row_h,
         "DIMENSION", 9, GRAY, bold=True)
add_text(slide, mat_x + label_w, mat_y + Inches(0.07), col_w_m, row_h,
         "BIG-BRAND OPERATORS", 9, AMBER, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, mat_x + label_w + col_w_m, mat_y + Inches(0.07), col_w_m, row_h,
         "SELF-SERVE PLATFORMS", 9, BLUE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, mat_x + label_w + col_w_m*2, mat_y + Inches(0.07), col_w_m, row_h,
         "KIX", 9, KIX_GREEN, bold=True, align=PP_ALIGN.CENTER)

rows = [
    # (dimension, operators, self-serve, KiX)
    ("Self-serve onboarding",        ("✗", "concierge"),         ("✓", "limited"),       ("✓", "5-minute setup")),
    ("Marginal cost per merchant",   ("✗", "$50K–$5M"),          ("≈", "$1K–$10K"),     ("✓", "≈ $0")),
    ("Pricing",                       ("✗", "custom / 6-fig"),    ("≈", "$500–$5K/mo"),  ("✓", "$49/month")),
    ("Languages",                     ("≈", "1–5 majors"),        ("≈", "EN-first"),      ("✓", "11 locales · CLDR")),
    ("SEA payment methods",          ("✗", "USD/EUR only"),      ("✗", "Stripe basics"), ("✓", "60 × 14 countries")),
    ("Cross-brand network",          ("✗", "one brand"),         ("✗", "one brand"),     ("✓", "marketplace")),
    ("Compliance coverage",          ("≈", "GDPR/CCPA"),         ("≈", "GDPR only"),     ("✓", "9 regions coded")),
    ("Time to launch a campaign",    ("✗", "6–12 weeks"),        ("≈", "1–2 weeks"),     ("✓", "5 minutes")),
    ("Target customer",              ("✗", "Fortune 500"),       ("≈", "mid-market"),    ("✓", "70M offline SMBs")),
]
y = mat_y + row_h
for i, (dim, op, ss, kx) in enumerate(rows):
    bg = NAVY_LIGHT if i % 2 == 0 else NAVY
    add_rect(slide, mat_x, y, mat_w, row_h, bg)
    add_text(slide, mat_x + Inches(0.15), y + Inches(0.08), label_w - Inches(0.15), row_h,
             dim, 9, WHITE, bold=True)
    # 3 cells with mark + tiny note
    for col_idx, (mark, note) in enumerate([op, ss, kx]):
        cx = mat_x + label_w + col_w_m * col_idx
        mark_color = KIX_GREEN if mark == "✓" else (RED_MUTED if mark == "✗" else AMBER)
        # mark
        tb = slide.shapes.add_textbox(cx, y + Inches(0.04), col_w_m, row_h)
        tf = tb.text_frame
        tf.margin_left = Emu(0); tf.margin_right = Emu(0)
        tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r_m = p.add_run(); r_m.text = mark + "  "
        r_m.font.name = "Inter"; r_m.font.size = Pt(11); r_m.font.bold = True
        r_m.font.color.rgb = mark_color
        r_n = p.add_run(); r_n.text = note
        r_n.font.name = "Inter"; r_n.font.size = Pt(8.5); r_n.font.bold = False
        r_n.font.color.rgb = GRAY_LIGHT
    y += row_h

# Bottom KiX tagline
tag_y = y + Inches(0.15)
add_text(slide, mat_x, tag_y, mat_w, Inches(0.4),
         "KiX is the only column that's green on every row that matters.",
         12, KIX_GREEN, bold=True, align=PP_ALIGN.CENTER)

add_footer(slide, "Differentiation Matrix · KiX = self-serve + cross-brand + SEA-native at marginal cost ≈ 0")


# ═══════════════════════════════════════════════════════════════════
# SLIDE 5 — System Architecture (updated v1 slide 2)
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
add_slide_header(slide, 5, "UNDER THE HOOD",
                 "An enterprise stack — built by AI in weeks.",
                 "One platform, three layers. Production-grade, multi-region, multi-locale, multi-currency.",
                 headline_size=28)

arch_cards = [
    ("FRONTEND",     "Vanilla JS · i18next · 25-component design system · Inter · WCAG AA · RTL-ready · ⌘K"),
    ("BACKEND",      "FastAPI · 85 routers · 1248 tests · PostgreSQL · Redis (HLL/Streams) · OAuth/JWT"),
    ("INTELLIGENCE", "LightGBM Smart Bidding · A/B testing · Multi-touch attribution · Soul behavior models"),
    ("MARKETPLACE",  "GSP Vickrey auction · PI-controller pacing · QS auto-compute · Viral K-factor compounding"),
    ("PAYMENT",      "60 methods × 14 countries · Stripe live/test/mock · GrabPay/OVO/PayNow/Alipay/WeChat"),
    ("INFRA",        "Docker CN/ID/SG · K8s · GeoDNS · PostGIS geofence · GDPR portability · CSV export"),
    ("I18N",         "11 locales · ICU · BCP-47 · CLDR plurals · LLM batch translator · $1.83/locale"),
    ("COMPLIANCE",   "9 regions: GDPR · PIPL · PDPA · CCPA · DPDP · LGPD · age gates + DPO flags"),
    ("OPS",          "Trinity 3T iteration (5 rounds · 102 → 0 issues) · LLM quota guard · Alpha cohort worker"),
]
ac_w = Inches(2.0); ac_h = Inches(1.0); ac_x0 = Inches(0.6); ac_y0 = Inches(3.15); gap = Inches(0.12)
for idx, (head, body) in enumerate(arch_cards):
    row, col = divmod(idx, 3)
    x = ac_x0 + col * (ac_w + gap)
    y = ac_y0 + row * (ac_h + gap)
    add_arch_card(slide, x, y, ac_w, ac_h, head, body)

# Screenshot right
ss_path = SCREENSHOT_DIR / "design-system.png"
if ss_path.exists():
    box_x = Inches(7.4); box_y = Inches(3.15); ss_w = Inches(5.3); ss_h = Inches(3.3)
    add_rect(slide, box_x, box_y, ss_w + Inches(0.2), ss_h + Inches(0.5),
             NAVY_LIGHT, NAVY_BORDER, 0.5)
    slide.shapes.add_picture(str(ss_path), box_x + Inches(0.1), box_y + Inches(0.1),
                             width=ss_w, height=ss_h)
    add_text(slide, box_x + Inches(0.1), box_y + ss_h + Inches(0.15), ss_w, Inches(0.4),
             "Merchant Portal · Enterprise component library · Reused everywhere",
             9, GRAY)

add_footer(slide, "Architecture · 285K LOC · 11 locales · 9 regions · 60 payment methods · 1248 tests")


# ═══════════════════════════════════════════════════════════════════
# SLIDE 6 — 90-Day Bedok Land Grab
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
add_slide_header(slide, 6, "90-DAY BEDOK LAND GRAB",
                 "Density before breadth. Lock-in before scale.",
                 "One neighbourhood. 100 merchants. 5,000 players. Irreversible network density.",
                 headline_size=28)

# Horizontal timeline: 4 milestone cards (Day 0, 30, 60, 90)
tl_y = Inches(3.4)
# baseline line
add_rect(slide, Inches(0.6), tl_y + Inches(0.9), Inches(12.13), Inches(0.04), KIX_GREEN)

milestones = [
    ("DAY 0",   "Launch",           "Alpha cohort onboarded · Bedok F&B targeting · S$500 credit each"),
    ("DAY 30",  "30 merchants live", "Game + Reward proven · First voucher redemptions · cohort #1 retention measured"),
    ("DAY 60",  "60 merchants",     "+ Register + Redeem · 2,000 players · cross-merchant traffic emerging"),
    ("DAY 90",  "100 merchants",    "+ 5,000 players · Return = repeat visit cohort · irreversible density"),
]
n = len(milestones)
total_w = Inches(12.13)
gap = Inches(0.25)
ms_w = Emu(int((total_w - gap * (n - 1)) / n))
ms_h = Inches(2.5)
x = Inches(0.6)
for i, (day, head, body) in enumerate(milestones):
    # Marker circle on baseline
    circ_d = Inches(0.4)
    cx = x + ms_w / 2 - circ_d / 2
    cy = tl_y + Inches(0.72)
    color = KIX_GREEN if i > 0 else BLUE
    marker = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, circ_d, circ_d)
    marker.fill.solid(); marker.fill.fore_color.rgb = color
    marker.line.color.rgb = NAVY; marker.line.width = Pt(2)
    marker.shadow.inherit = False

    # Day label above
    add_text(slide, x, tl_y, ms_w, Inches(0.4),
             day, 14, color, bold=True, align=PP_ALIGN.CENTER)

    # Card below baseline
    card_y = tl_y + Inches(1.25)
    add_round_rect(slide, x, card_y, ms_w, ms_h, NAVY_LIGHT, color, 1.0, radius=0.06)
    pad = Inches(0.18)
    add_text(slide, x + pad, card_y + Inches(0.2), ms_w - pad*2, Inches(0.4),
             head, 13, WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x + pad, card_y + Inches(0.7), ms_w - pad*2, ms_h - Inches(0.85),
             body, 10, GRAY_LIGHT, align=PP_ALIGN.CENTER)
    x += ms_w + gap

# Footer callout
fc_y = Inches(7.05)
add_text(slide, Inches(0.6), fc_y - Inches(0.05), Inches(12.13), Inches(0.4),
         "Day 90 = irreversible network density.",
         13, KIX_GREEN, bold=True, align=PP_ALIGN.CENTER)

# Small inline alpha screenshot strip (right corner of headline area)
ss_path = SCREENSHOT_DIR / "alpha.png"
if ss_path.exists():
    box_x = Inches(9.4); box_y = Inches(0.4); ss_w = Inches(3.3); ss_h = Inches(2.0)
    add_rect(slide, box_x, box_y, ss_w + Inches(0.15), ss_h + Inches(0.4),
             NAVY_LIGHT, NAVY_BORDER, 0.5)
    slide.shapes.add_picture(str(ss_path), box_x + Inches(0.075), box_y + Inches(0.075),
                             width=ss_w, height=ss_h)
    add_text(slide, box_x + Inches(0.075), box_y + ss_h + Inches(0.1), ss_w, Inches(0.3),
             "Alpha · invite-only Bedok F&B pilot · 90 days free + S$500",
             7, GRAY)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 7 — The Crystal Clear Call
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
fill_bg(slide, NAVY)

add_logo(slide, Inches(0.6), Inches(0.4))
add_text(slide, Inches(0.6), Inches(0.95), Inches(3), Inches(0.25),
         f"SLIDE 7 / {TOTAL_SLIDES}", 8, GRAY_DARK, bold=True)
add_pill(slide, Inches(0.6), Inches(1.3), "THE ASK", width=2.0)

# Three big lines centered
lines = [
    ("KiX is the productized McDonald's Monopoly.", WHITE),
    ("For 70M offline merchants.",                  KIX_GREEN),
    ("Available in 5 minutes for $49/month.",        WHITE),
]
ly = Inches(2.0)
line_h = Inches(0.95)
for i, (txt, col) in enumerate(lines):
    add_text(slide, Inches(0.6), ly + i * line_h, Inches(12.13), line_h,
             txt, 36, col, bold=True, align=PP_ALIGN.CENTER)

# Divider
add_rect(slide, Inches(5.5), Inches(5.2), Inches(2.33), Inches(0.04), KIX_GREEN)

# Founder mantra block
mantra_y = Inches(5.5)
add_text(slide, Inches(0.6), mantra_y, Inches(12.13), Inches(0.35),
         "FOUNDER MANTRA",
         10, KIX_GREEN, bold=True, align=PP_ALIGN.CENTER)

# 4-up mantra row
mantras = [
    ("Speed",   "> perfection"),
    ("Density", "> breadth"),
    ("Lock-in", "> revenue"),
    ("Today",   "> tomorrow"),
]
mn = len(mantras); mgap = Inches(0.15)
m_total = Inches(10.0)
m_w = Emu(int((m_total - mgap * (mn - 1)) / mn))
m_h = Inches(0.7)
m_y = Inches(5.95)
m_x = Inches(0.6) + (Inches(12.13) - m_total) / 2  # center
for i, (a, b) in enumerate(mantras):
    x = m_x + i * (m_w + mgap)
    add_round_rect(slide, x, m_y, m_w, m_h, NAVY_LIGHT, KIX_GREEN, 0.75, radius=0.2)
    tb = slide.shapes.add_textbox(x, m_y + Inches(0.18), m_w, Inches(0.4))
    tf = tb.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r1 = p.add_run(); r1.text = a + " "
    r1.font.name = "Inter"; r1.font.size = Pt(14); r1.font.bold = True
    r1.font.color.rgb = KIX_GREEN
    r2 = p.add_run(); r2.text = b
    r2.font.name = "Inter"; r2.font.size = Pt(14); r2.font.bold = False
    r2.font.color.rgb = WHITE

# Contact
add_text(slide, Inches(0.6), Inches(6.9), Inches(12.13), Inches(0.4),
         "michael@mozat.com",
         14, GRAY_LIGHT, bold=True, align=PP_ALIGN.CENTER)

add_footer(slide, "KiX · The productized McDonald's Monopoly · 2026-05-30")


# ───── save ─────
prs.save(str(OUT_PATH))
print(f"✓ Saved: {OUT_PATH}")
print(f"  Size: {OUT_PATH.stat().st_size / 1024:.1f} KB")
print(f"  Slides: {len(prs.slides)}")
